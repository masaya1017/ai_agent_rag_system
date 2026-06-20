"""
RAG MCP Server - ローカルドキュメント検索エンジン
Claude Code から search_documents / ingest_document ツールとして呼び出す
"""
import os
import sys
import hashlib
from pathlib import Path
from typing import Optional

# ChromaDB と sentence-transformers は起動時に読み込む
import chromadb
from sentence_transformers import SentenceTransformer
import pypdf
import docx
from pptx import Presentation
from mcp.server.fastmcp import FastMCP

# ---- 設定 ----------------------------------------------------------------
DB_PATH = Path(__file__).parent / "db"
MODEL_NAME = "intfloat/multilingual-e5-large"  # 日本語対応の多言語埋め込みモデル
CHUNK_SIZE = 500   # チャンクあたりの文字数
CHUNK_OVERLAP = 50  # チャンク間のオーバーラップ文字数

# ---- 初期化 --------------------------------------------------------------
print(f"[RAG] 埋め込みモデル読み込み中: {MODEL_NAME}", file=sys.stderr)
embedding_model = SentenceTransformer(MODEL_NAME)
print("[RAG] モデル読み込み完了", file=sys.stderr)

chroma_client = chromadb.PersistentClient(path=str(DB_PATH))
collection = chroma_client.get_or_create_collection(
    name="documents",
    metadata={"hnsw:space": "cosine"},
)

mcp = FastMCP("RAG System")

# ---- ユーティリティ -------------------------------------------------------

def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    """テキストをオーバーラップ付きでチャンク分割する"""
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunks.append(text[start:end])
        start += chunk_size - overlap
    return [c for c in chunks if c.strip()]


def extract_text_from_pdf(path: Path) -> str:
    reader = pypdf.PdfReader(str(path))
    pages = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            pages.append(text)
    return "\n".join(pages)


def extract_text_from_docx(path: Path) -> str:
    doc = docx.Document(str(path))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def extract_text_from_pptx(path: Path) -> str:
    prs = Presentation(str(path))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        texts.append(line)
    return "\n".join(texts)


def extract_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return extract_text_from_pdf(path)
    elif suffix in (".docx", ".doc"):
        return extract_text_from_docx(path)
    elif suffix in (".pptx", ".ppt"):
        return extract_text_from_pptx(path)
    elif suffix in (".txt", ".md"):
        return path.read_text(encoding="utf-8", errors="replace")
    else:
        raise ValueError(f"未対応のファイル形式: {suffix}")


def embed(texts: list[str]) -> list[list[float]]:
    # multilingual-e5 は "query: " / "passage: " プレフィックスが必要
    prefixed = [f"passage: {t}" for t in texts]
    return embedding_model.encode(prefixed, normalize_embeddings=True).tolist()


def file_id(path: Path) -> str:
    return hashlib.md5(str(path).encode()).hexdigest()

# ---- MCPツール ------------------------------------------------------------

@mcp.tool()
def ingest_document(file_path: str) -> str:
    """
    ドキュメント（PDF/Word/テキスト）をベクターDBに取り込む。

    Args:
        file_path: 取り込むファイルの絶対パスまたは相対パス
    """
    path = Path(file_path).expanduser().resolve()
    if not path.exists():
        return f"エラー: ファイルが見つかりません: {path}"

    try:
        text = extract_text(path)
    except ValueError as e:
        return f"エラー: {e}"
    except Exception as e:
        return f"テキスト抽出エラー: {e}"

    if not text.strip():
        return f"警告: {path.name} からテキストを抽出できませんでした"

    chunks = chunk_text(text)
    doc_id = file_id(path)

    # 既存チャンクを削除（再取り込み対応）
    existing = collection.get(where={"source_file": str(path)})
    if existing["ids"]:
        collection.delete(ids=existing["ids"])

    # チャンクを埋め込んでDBに保存
    chunk_ids = [f"{doc_id}_{i}" for i in range(len(chunks))]
    embeddings = embed(chunks)
    metadatas = [
        {"source_file": str(path), "filename": path.name, "chunk_index": i}
        for i in range(len(chunks))
    ]

    collection.add(
        ids=chunk_ids,
        embeddings=embeddings,
        documents=chunks,
        metadatas=metadatas,
    )

    return (
        f"取り込み完了: {path.name}\n"
        f"  チャンク数: {len(chunks)}\n"
        f"  文字数: {len(text)}"
    )


@mcp.tool()
def ingest_directory(directory_path: str, recursive: bool = True) -> str:
    """
    ディレクトリ内のすべてのPDF/Word/テキストファイルを一括取り込みする。

    Args:
        directory_path: 取り込むディレクトリのパス
        recursive: サブディレクトリも含めるか（デフォルト: True）
    """
    dir_path = Path(directory_path).expanduser().resolve()
    if not dir_path.is_dir():
        return f"エラー: ディレクトリが見つかりません: {dir_path}"

    pattern = "**/*" if recursive else "*"
    extensions = {".pdf", ".docx", ".doc", ".pptx", ".ppt", ".txt", ".md"}
    files = [f for f in dir_path.glob(pattern) if f.is_file() and f.suffix.lower() in extensions]

    if not files:
        return f"対象ファイルが見つかりません: {dir_path}"

    results = []
    for f in files:
        result = ingest_document(str(f))
        results.append(f"[{f.name}] {result.split(chr(10))[0]}")

    return f"一括取り込み完了: {len(files)} ファイル\n" + "\n".join(results)


@mcp.tool()
def search_documents(query: str, n_results: int = 5, source_filter: Optional[str] = None) -> str:
    """
    取り込み済みドキュメントから質問に関連する情報を検索する。

    Args:
        query: 検索クエリ（日本語可）
        n_results: 返す結果の最大数（デフォルト: 5）
        source_filter: ファイル名でフィルタリング（省略可）
    """
    total = collection.count()
    if total == 0:
        return "ドキュメントがまだ取り込まれていません。ingest_document を使って取り込んでください。"

    # クエリを埋め込み
    query_embedding = embedding_model.encode(
        f"query: {query}", normalize_embeddings=True
    ).tolist()

    where = {}
    if source_filter:
        where["filename"] = {"$contains": source_filter}

    results = collection.query(
        query_embeddings=[query_embedding],
        n_results=min(n_results, total),
        where=where if where else None,
        include=["documents", "metadatas", "distances"],
    )

    if not results["ids"][0]:
        return "関連するドキュメントが見つかりませんでした。"

    output_parts = [f"「{query}」の検索結果 ({len(results['ids'][0])} 件):\n"]
    for i, (doc, meta, dist) in enumerate(
        zip(results["documents"][0], results["metadatas"][0], results["distances"][0])
    ):
        similarity = 1 - dist  # コサイン距離 → 類似度
        output_parts.append(
            f"--- [{i+1}] {meta['filename']} (類似度: {similarity:.2f}) ---\n{doc}\n"
        )

    return "\n".join(output_parts)


@mcp.tool()
def list_documents() -> str:
    """取り込み済みのドキュメント一覧と統計を表示する。"""
    total_chunks = collection.count()
    if total_chunks == 0:
        return "ドキュメントがまだ取り込まれていません。"

    all_meta = collection.get(include=["metadatas"])["metadatas"]
    files: dict[str, int] = {}
    for meta in all_meta:
        fn = meta.get("filename", "unknown")
        files[fn] = files.get(fn, 0) + 1

    lines = [f"取り込み済みドキュメント: {len(files)} ファイル / {total_chunks} チャンク\n"]
    for filename, chunk_count in sorted(files.items()):
        lines.append(f"  • {filename} ({chunk_count} チャンク)")

    return "\n".join(lines)


@mcp.tool()
def delete_document(filename: str) -> str:
    """
    指定したファイル名のドキュメントをベクターDBから削除する。

    Args:
        filename: 削除するファイル名（例: document.pdf）
    """
    existing = collection.get(where={"filename": {"$eq": filename}})
    if not existing["ids"]:
        return f"ドキュメントが見つかりません: {filename}"

    collection.delete(ids=existing["ids"])
    return f"削除完了: {filename} ({len(existing['ids'])} チャンク)"


if __name__ == "__main__":
    print("[RAG] MCPサーバー起動", file=sys.stderr)
    mcp.run()
